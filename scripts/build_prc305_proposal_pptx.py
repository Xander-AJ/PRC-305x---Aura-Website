#!/usr/bin/env python3
"""Generate PRC 305 (Jan 2026) proposal deck per course guidelines — Aura Energy project."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Brand-aligned accent (navy + gold) for headings
NAVY = RGBColor(6, 13, 24)
GOLD = RGBColor(227, 199, 111)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(80, 90, 100)


def set_slide_title(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT


def add_bullets(shape, lines: list[str], size: int = 18) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
        p.level = 0


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Gold bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.2))
    t = title_box.text_frame
    t.word_wrap = True
    p = t.paragraphs[0]
    p.text = "PRC 305 — Web Design & CMS for PR"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p2 = t.add_paragraph()
    p2.text = "Project Proposal"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.space_before = Pt(12)
    p3 = t.add_paragraph()
    p3.text = "Aura Energy — Premium Functional Beverages (East Africa)"
    p3.font.size = Pt(22)
    p3.font.color.rgb = GOLD
    p3.space_before = Pt(18)
    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11), Inches(1.5))
    mt = meta.text_frame
    mt.clear()
    mlines = [
        "[Your Name]  ·  [Student ID]",
        "School of Communication  ·  January–April 2026",
        "Planned stack: WordPress, Hello Elementor, Elementor, MetForm",
    ]
    for i, line in enumerate(mlines):
        mp = mt.paragraphs[0] if i == 0 else mt.add_paragraph()
        mp.text = line
        mp.font.size = Pt(14)
        mp.font.color.rgb = GRAY

    # --- Slide 2: Domain & background ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide.shapes.title, "Selected domain & background")
    body = slide.placeholders[1]
    add_bullets(
        body,
        [
            "Domain: Consumer beverage / brand communications — premium functional & energy drinks positioned for Kenya and East Africa.",
            "Context: Growing demand for better-for-you, lower-sugar options among urban professionals, creatives, and shift workers who expect credible branding and mobile-first discovery.",
            "Organizational concept: “Aura Energy” — a focused product family (Focus, Flow, Night) plus packaging and carry concepts, presented as a realistic regional launch story.",
            "Why it matters: Bridges PR, brand narrative, and digital touchpoints (web as primary credibility + conversion surface for stockists, events, and press).",
        ],
        17,
    )

    # --- Slide 3: Executive summary / company profile ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide.shapes.title, "Executive summary / company profile")
    body = slide.placeholders[1]
    add_bullets(
        body,
        [
            "Aura Energy crafts functional beverages for sustained focus, hydration, and calmer wind-down — without relying on neon-sugar clichés.",
            "Brand promise: honest formulation, East African sensibility (pace, climate, culture), and visual identity built on deep navy, jewel teal, and champagne gold.",
            "Digital presence goal: a polished multi-page site that mirrors premium packaging, explains the lineup, and makes partnership inquiries effortless.",
            "This proposal scopes the WordPress website that will support launch communications, B2B outreach, and public trust.",
        ],
        17,
    )

    # --- Slide 4: Objectives & justification ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide.shapes.title, "Project objectives & justification")
    body = slide.placeholders[1]
    add_bullets(
        body,
        [
            "Objectives: (1) Publish a responsive 4-page marketing site (Home, About, Products, Contact). (2) Implement a validated contact form routed to brand/admin email. (3) Integrate ≥2 social channels. (4) Maintain consistent UI (typography, color, spacing) and optional gallery.",
            "Justification: Meets PRC 305 mandatory requirements while modeling a real communications use case — stakeholders judge beverage brands heavily on web quality, speed, and clarity.",
            "Success: Site is usable on phone/tablet/desktop, form works, navigation is obvious, and content reads as professional (not template filler).",
        ],
        16,
    )

    # --- Slide 5: Target audience & UX ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide.shapes.title, "Target audience & key UX considerations")
    body = slide.placeholders[1]
    add_bullets(
        body,
        [
            "Primary audience: Young professionals & creatives (25–40) in Nairobi and regional cities; secondary: retail buyers, event organizers, lifestyle press.",
            "UX priorities: Fast scan-ability, thumb-friendly CTAs, high-contrast readable type, credible product storytelling, minimal friction to “Contact / Partner”.",
            "Mobile-first: Single-column stacks, tappable social icons, form fields with clear labels and validation feedback (MetForm).",
            "Trust signals: Consistent Aura visuals, real-sounding copy, map/studio block on Contact, professional footer and legal line.",
        ],
        16,
    )

    # --- Slide 6: Sitemap ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide.shapes.title, "Sitemap & information architecture")
    body = slide.placeholders[1]
    add_bullets(
        body,
        [
            "Home — Hero + brand promise; About teaser; product highlights; optional gallery; CTA to Contact; footer + social.",
            "About — Company profile, brand story, visual anchor (hero image).",
            "Products / Services — Lineup (Focus / Flow / Night), hero SKU detail, carry/activation merch.",
            "Contact — Intro copy, MetForm (name, email, message + validation + success state), studio hours, embedded map, social.",
            "Global: Primary menu (4 items), static front page = Home, Elementor Site Settings for global colors & fonts.",
        ],
        15,
    )

    # --- Slide 7: Wireframes / layout sketches ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide.shapes.title, "Preliminary wireframes & layout approach")
    body = slide.placeholders[1]
    add_bullets(
        body,
        [
            "Home (Z-pattern): Top = logo + nav; hero = headline + subcopy + dual CTAs + product visual; mid = two-column About teaser; three-column SKU cards; image grid; bottom CTA band + footer.",
            "About / Products: Single-column narrative + full-width imagery; Products adds split rows (image | specs list).",
            "Contact: Narrow content width for readability; form block + map below; repeated footer pattern for consistency.",
            "Sketches: [Insert hand-drawn or Figma frames here before presentation] — align sections to Elementor sections/columns for responsive breakpoints.",
        ],
        15,
    )

    out = Path(__file__).resolve().parent.parent / "PRC305-Aura-Project-Proposal.pptx"
    prs.save(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
