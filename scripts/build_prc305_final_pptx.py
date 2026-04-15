#!/usr/bin/env python3
"""PRC 305 (Jan 2026) FINAL presentation — Aura Energy website. Submission-ready deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(6, 13, 24)
GOLD = RGBColor(227, 199, 111)
GRAY = RGBColor(75, 85, 99)
DARK = RGBColor(40, 45, 55)

STUDENT_NAME = "John Alexander Kamau"
STUDENT_ID = "20-0471"
PROJECT = "Aura Energy — Premium Functional Beverages (East Africa)"
# Update if deployed publicly:
LIVE_URL = "http://localhost/wordpress/"
REPO_NOTE = "GitHub: [paste your repository URL here]"


def title_style(paragraph, size: int, bold: bool = False, color=GRAY):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(11.8), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PRC 305: Web Design & Content Management Systems for PR"
    title_style(p, 18, False, GRAY)
    p2 = tf.add_paragraph()
    p2.text = "Final Project Presentation"
    title_style(p2, 38, True, NAVY)
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = PROJECT
    title_style(p3, 22, False, GOLD)
    p3.space_before = Pt(16)

    meta = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(11.5), Inches(2))
    m = meta.text_frame
    m.clear()
    lines = [
        STUDENT_NAME,
        f"Student No. {STUDENT_ID}",
        "School of Communication  ·  January–April 2026",
        "Submission date: 1 April 2026 (per course outline)",
    ]
    for i, line in enumerate(lines):
        mp = m.paragraphs[0] if i == 0 else m.add_paragraph()
        mp.text = line
        title_style(mp, 15, i < 2, DARK if i < 2 else GRAY)
        mp.space_after = Pt(6)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], size: int = 17) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(size)
        para.font.color.rgb = GRAY
        para.space_after = Pt(8)
        para.level = 0


def add_section_divider(prs: Presentation, label: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.5), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 210, 220)
        p2.space_before = Pt(14)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "Agenda",
        [
            "Introduction & background",
            "Objectives & problem statement",
            "Design process & UI/UX",
            "Technology stack & development approach",
            "Deliverables & site walkthrough",
            "Challenges, solutions & lessons",
            "Recommendations, conclusion & Q&A",
        ],
        19,
    )

    add_bullet_slide(
        prs,
        "Project introduction & background",
        [
            "Aura Energy is a premium functional beverage concept aimed at young professionals, creators, and shift workers in Kenya and East Africa.",
            "The website supports brand credibility, explains the product lineup (Focus · Flow · Night), and channels wholesale, events, and press inquiries.",
            "The project applies the full web lifecycle: planning, design in Elementor, implementation in WordPress, testing, and documentation.",
            "Aligned with PRC 305: realistic scenario, CMS-based multi-page site, responsive layout, and professional presentation.",
        ],
    )

    add_bullet_slide(
        prs,
        "Problem statement",
        [
            "Functional drink brands are judged quickly online: weak sites reduce trust with buyers, media, and end consumers.",
            "Stakeholders need clear product story, social proof, frictionless contact, and mobile-first access in a high-mobile region.",
            "A generic template site fails to communicate premium positioning or support B2B outreach.",
            "This project addresses that gap with a cohesive brand site built on WordPress best practices.",
        ],
    )

    add_bullet_slide(
        prs,
        "Objectives",
        [
            "Deliver a fully functional WordPress site with Home, About, Products, and Contact pages.",
            "Implement a working contact form with validation and success feedback (MetForm).",
            "Integrate at least two social platforms (Instagram, Facebook, YouTube) with clear labels/placeholders.",
            "Ensure responsive behavior across mobile, tablet, and desktop (Elementor responsive controls).",
            "Maintain visual consistency: navy / teal / champagne gold palette, typography, spacing, and Cloudinary imagery.",
        ],
    )

    add_section_divider(prs, "Design & UX", "Brand-led, user-centered")

    add_bullet_slide(
        prs,
        "Design process & UI/UX",
        [
            "Brand foundation: deep navy (#060d18), jewel teal gradients, champagne gold accents — premium, not “neon energy” cliché.",
            "Typography: Outfit + DM Sans (Elementor Site Settings / Kit) for clear hierarchy and readability.",
            "Layout: hero with dual CTAs, teaser + product grid, optional gallery, strong footer with social icons.",
            "UX choices: obvious primary navigation, scannable copy, tappable targets on mobile, map + studio block on Contact for trust.",
            "Imagery: Cloudinary-hosted brand assets (q_auto, f_auto) for performance.",
        ],
    )

    add_bullet_slide(
        prs,
        "Information architecture",
        [
            "Home — Static front page: hero, about teaser, three SKU highlights, design gallery, CTA band, footer.",
            "About — Brand story, audience, tone, and visual anchor.",
            "Products — Lineup narrative, hero SKU (Aura Original), carry / activation merch.",
            "Contact — Intro, MetForm, studio details, embedded map, footer + social (mirrors site chrome).",
            "Global: Primary menu assigned to Hello Elementor header location; reading settings = static homepage.",
        ],
    )

    add_section_divider(prs, "Technology", "Stack & workflow")

    add_bullet_slide(
        prs,
        "Technology stack",
        [
            "CMS: WordPress 6.x on XAMPP (local) — exportable to any LAMP host for live URL submission.",
            "Theme: Hello Elementor (lightweight, Elementor-native).",
            "Page builder: Elementor (free) — sections, columns, widgets; Site Settings for global colors & fonts.",
            "Forms: MetForm — name, email, message; validation; admin notification & entries.",
            "Assets: Cloudinary URLs embedded in Image widgets; optional Elementor JSON templates in repo for reproducibility.",
        ],
    )

    add_bullet_slide(
        prs,
        "Development approach",
        [
            "Structured setup script (`aura-prc305-setup.php`) to sync Elementor JSON, kit, pages, menu, and MetForm — reduces drift between export files and database.",
            "Elementor library JSON (version 0.4 pattern) stored under `elementor-templates/` for course submission alongside the running site.",
            "Iterative refinement of copy and palette for a more “launch-ready” brand voice.",
            "Testing: front-end checks per breakpoint; form submission to admin email; link checks for menu and CTAs.",
        ],
    )

    add_section_divider(prs, "Demonstration", "Live site")

    add_bullet_slide(
        prs,
        "Live demonstration — what to show",
        [
            f"Open the site: {LIVE_URL}",
            "Walk through: Home → About → Products → Contact; show responsive preview or resize browser.",
            "Demonstrate Contact form: empty submit (validation), valid submit (success message / entry).",
            "Point out social icons (footer / contact) and gallery or product imagery.",
            "If hosted publicly for submission: replace localhost URL on this slide with production link.",
        ],
    )

    add_bullet_slide(
        prs,
        "Submission materials (checklist)",
        [
            f"Live website URL: {LIVE_URL} (update when deployed)",
            f"Source / project evidence: {REPO_NOTE}",
            "This presentation (final slides) + any written documentation required by instructor.",
            "Elementor template JSON files in project folder for import/audit if requested.",
        ],
    )

    add_section_divider(prs, "Reflection", "Challenges & outcomes")

    add_bullet_slide(
        prs,
        "Challenges & solutions",
        [
            "Elementor data format: full template JSON must split into `content` array for post meta — scripted save avoided front-end fatals.",
            "MetForm default templates included reCAPTCHA without keys — removed widget + enabled admin email notifications.",
            "Menu order & static homepage — set via script/Customizer so IA matches assignment rubric.",
            "Plugin noise: My Sticky Menu deactivated where it conflicted with a clean header experience.",
        ],
    )

    add_bullet_slide(
        prs,
        "Recommendations & conclusion",
        [
            "Deploy to a public host with HTTPS before final deadline; configure SMTP for reliable form email.",
            "Replace placeholder social URLs with real brand handles when accounts go live.",
            "Consider lightweight caching and image CDN (Cloudinary already helps) for production performance.",
            "Conclusion: The Aura site meets PRC 305 requirements — multi-page CMS site, validated form, social integration, responsive design, and consistent visual system.",
        ],
    )

    add_bullet_slide(
        prs,
        "Alignment with assessment (50% of course grade)",
        [
            "Functionality & completeness (30%) — Four core pages, working MetForm with validation, social links, optional gallery.",
            "Design quality & responsiveness (25%) — Cohesive Aura palette, typography, imagery, Elementor responsive tuning.",
            "Web structure (20%) — Clear IA, primary menu, static homepage, logical CTAs.",
            "Presentation & documentation (15%) — This deck + demo + any required written docs.",
            "Deployment & live demonstration (10%) — Public URL when deployed; local demo acceptable only if instructor allows.",
        ],
        15,
    )

    add_bullet_slide(
        prs,
        "Thank you — Q&A",
        [
            f"{STUDENT_NAME}  ·  Student No. {STUDENT_ID}",
            "I welcome your questions.",
            "Aura Energy — student web development project (PRC 305).",
        ],
        22,
    )

    out = Path(__file__).resolve().parent.parent / "PRC305-Aura-Final-Presentation.pptx"
    prs.save(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
